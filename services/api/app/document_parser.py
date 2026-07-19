from dataclasses import dataclass, field
from pathlib import Path

from .models import SourceType


@dataclass
class ParsedChunk:
    text: str
    heading: str | None = None
    page: int | None = None


@dataclass
class ParsedDocument:
    source_type: SourceType
    chunks: list[ParsedChunk]
    parser: str
    warnings: list[str] = field(default_factory=list)

    @property
    def text(self) -> str:
        return "\n\n".join(chunk.text for chunk in self.chunks)


LLM_TASKS = [
    "extract_company_profile",
    "extract_founders",
    "extract_traction_metrics",
    "extract_funding_terms",
    "link_claims_to_evidence",
    "flag_missing_information",
]


def parse_document(filename: str, content: bytes) -> ParsedDocument:
    ext = Path(filename).suffix.lower()
    if ext in {".txt", ".md"}:
        return _parse_text(content, _source_type(ext))
    if ext == ".pdf":
        return _parse_pdf(content)
    if ext == ".pptx":
        return _parse_pptx(content)
    if ext == ".docx":
        return _parse_docx(content)
    if ext in {".xlsx", ".xls"}:
        return _parse_spreadsheet(content, ext)
    if ext in {".png", ".jpg", ".jpeg", ".webp"}:
        return _parse_image(content, ext)
    return _parse_text(content, SourceType.document, warning=f"Unsupported extension: {ext or 'none'}")


def _source_type(ext: str) -> SourceType:
    if ext in {".pdf", ".pptx"}:
        return SourceType.pitch_deck
    if ext == ".docx":
        return SourceType.founder_questionnaire
    return SourceType.document


def _parse_text(content: bytes, source_type: SourceType, warning: str | None = None) -> ParsedDocument:
    text = content.decode("utf-8", errors="ignore").strip()
    chunks = _chunk_text(text or "No readable text extracted.", "Document")
    warnings = [warning] if warning else []
    return ParsedDocument(source_type=source_type, chunks=chunks, parser="text", warnings=warnings)


def _parse_pdf(content: bytes) -> ParsedDocument:
    try:
        from pypdf import PdfReader
        from io import BytesIO

        reader = PdfReader(BytesIO(content))
        chunks = [
            ParsedChunk(
                page=index,
                heading=f"Page {index}",
                text=(page.extract_text() or "").strip() or "No readable text on page.",
            )
            for index, page in enumerate(reader.pages, start=1)
        ]
        return ParsedDocument(SourceType.pitch_deck, chunks, "pypdf")
    except Exception as exc:
        fallback = _parse_text(content, SourceType.pitch_deck)
        fallback.parser = "pdf_fallback"
        fallback.warnings.append(f"PDF parser unavailable or failed: {type(exc).__name__}")
        return fallback


def _parse_pptx(content: bytes) -> ParsedDocument:
    try:
        from io import BytesIO
        from pptx import Presentation

        deck = Presentation(BytesIO(content))
        chunks: list[ParsedChunk] = []
        for index, slide in enumerate(deck.slides, start=1):
            text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            chunks.append(
                ParsedChunk(
                    page=index,
                    heading=f"Slide {index}",
                    text=text or "No readable text on slide.",
                )
            )
        return ParsedDocument(SourceType.pitch_deck, chunks, "python-pptx")
    except Exception as exc:
        fallback = _parse_text(content, SourceType.pitch_deck)
        fallback.parser = "pptx_fallback"
        fallback.warnings.append(f"PPTX parser unavailable or failed: {type(exc).__name__}")
        return fallback


def _parse_docx(content: bytes) -> ParsedDocument:
    try:
        from io import BytesIO
        from docx import Document

        doc = Document(BytesIO(content))
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
        return ParsedDocument(SourceType.founder_questionnaire, _chunk_text(text, "Document"), "python-docx")
    except Exception as exc:
        fallback = _parse_text(content, SourceType.founder_questionnaire)
        fallback.parser = "docx_fallback"
        fallback.warnings.append(f"DOCX parser unavailable or failed: {type(exc).__name__}")
        return fallback


def _parse_spreadsheet(content: bytes, ext: str) -> ParsedDocument:
    try:
        from io import BytesIO
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        chunks: list[ParsedChunk] = []
        for sheet in workbook.worksheets:
            rows = []
            for values in sheet.iter_rows(values_only=True):
                cells = [str(value).strip() for value in values if value is not None and str(value).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            chunks.append(ParsedChunk(heading=f"Sheet: {sheet.title}", text="\n".join(rows) or "Empty sheet."))
        return ParsedDocument(SourceType.financial_model, chunks, "openpyxl")
    except Exception as exc:
        fallback = _parse_text(content, SourceType.financial_model)
        fallback.parser = "spreadsheet_fallback"
        fallback.warnings.append(f"Spreadsheet parser unavailable or failed: {type(exc).__name__}")
        return fallback


def _parse_image(content: bytes, ext: str) -> ParsedDocument:
    warnings: list[str] = []
    try:
        from io import BytesIO
        from PIL import Image

        image = Image.open(BytesIO(content))
        metadata = f"Image dimensions: {image.width}x{image.height}; format: {image.format or ext.lstrip('.')}."
        try:
            import pytesseract
            text = pytesseract.image_to_string(image).strip()
            if text:
                return ParsedDocument(SourceType.founder_questionnaire, _chunk_text(text, "Image OCR"), "pytesseract")
            warnings.append("OCR ran but found no readable text.")
        except Exception:
            warnings.append("OCR is optional; install Tesseract to extract image text locally.")
        return ParsedDocument(SourceType.founder_questionnaire, [ParsedChunk(heading="Image", text=metadata)], "pillow", warnings)
    except Exception as exc:
        return ParsedDocument(
            SourceType.founder_questionnaire,
            [ParsedChunk(heading="Image", text="Image received; no readable text extracted.")],
            "image_fallback",
            [f"Image parser unavailable or failed: {type(exc).__name__}"],
        )


def _chunk_text(text: str, heading: str) -> list[ParsedChunk]:
    size = 1800
    parts = [text[index : index + size].strip() for index in range(0, len(text), size)]
    return [
        ParsedChunk(page=index, heading=f"{heading} chunk {index}", text=part)
        for index, part in enumerate(parts or ["No readable text extracted."], start=1)
    ]
